"""pptx資料生成モジュール - 3C分析レポート"""

import io
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from config import PptxStyle
from models import ResearchResult


# --- フォント設定 ---

def _get_jp_font():
    """利用可能な日本語フォントを探す"""
    jp_fonts = [
        "Noto Sans JP", "Noto Sans CJK JP", "Hiragino Sans",
        "Hiragino Kaku Gothic ProN", "Yu Gothic", "Meiryo",
        "IPAGothic", "IPAPGothic",
    ]
    available = {f.name for f in fm.fontManager.ttflist}
    for font in jp_fonts:
        if font in available:
            return font
    return "sans-serif"


JP_FONT = _get_jp_font()
plt.rcParams["font.family"] = JP_FONT
plt.rcParams["axes.unicode_minus"] = False


# --- ヘルパー関数 ---

def _hex_to_rgb(hex_str: str) -> RGBColor:
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _add_text(tf, text, size=12, bold=False, color=None, alignment=None):
    """テキストフレームに段落を追加"""
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    if tf.paragraphs[0].text:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = _hex_to_rgb(color)
    try:
        run.font.name = PptxStyle.FONT_TITLE
    except Exception:
        run.font.name = PptxStyle.FONT_FALLBACK
    if alignment:
        p.alignment = alignment
    return p


def _set_shape_bg(shape, hex_color):
    """図形の背景色を設定"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _add_blank_slide(prs):
    """ブランクスライドを追加し、デフォルト要素を全て削除して返す"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # デフォルトのプレースホルダー図形（緑/ピンクの四角等）を全削除
    sp_tree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        sp = ph._element
        sp_tree.remove(sp)
    # レイアウトから継承された非プレースホルダー要素も削除
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
             "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
             "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}
    # sp要素（Shape）でidx属性を持つものを除去
    for sp in list(sp_tree):
        if sp.tag.endswith("}sp"):
            nvSpPr = sp.find(".//{%s}nvSpPr" % nsmap["p"])
            if nvSpPr is not None:
                nvPr = nvSpPr.find("{%s}nvPr" % nsmap["p"])
                if nvPr is not None and nvPr.find("{%s}ph" % nsmap["p"]) is not None:
                    sp_tree.remove(sp)
    return slide


def _draw_timeline_shapes(slide, events):
    """タイムライン（python-pptx図形で構築 - PPTX上で編集可能）"""
    n = len(events)
    if n == 0:
        return

    # タイムライン横線
    line_y = Inches(2.2)
    line_left = Inches(0.8)
    line_width = Inches(10.4)
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, line_left, line_y, line_width, Pt(3)
    )
    _set_shape_bg(line_shape, PptxStyle.ACCENT)
    line_shape.line.fill.background()

    # 各イベントの配置
    spacing = 10.4 / max(n, 1)
    for i, event in enumerate(events):
        x_center = 0.8 + spacing * (i + 0.5)

        # 丸マーカー
        dot_size = 0.2
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x_center - dot_size / 2), Inches(2.2 - dot_size * 0.4),
            Inches(dot_size), Inches(dot_size),
        )
        _set_shape_bg(dot, PptxStyle.ACCENT)
        dot.line.fill.background()

        # テキスト（上下交互配置）
        desc = event.description[:40] + "…" if len(event.description) > 40 else event.description
        label_text = f"{event.year}\n{desc}"
        tb_width = min(spacing + 0.1, 1.8)

        if i % 2 == 0:
            # 上に配置
            tb_y = 1.0
        else:
            # 下に配置
            tb_y = 2.6

        txBox = slide.shapes.add_textbox(
            Inches(x_center - tb_width / 2), Inches(tb_y),
            Inches(tb_width), Inches(1.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _add_text(tf, label_text, size=8, color=PptxStyle.TEXT_DARK, alignment=PP_ALIGN.CENTER)


# --- チャート生成（matplotlib → 画像バイト） ---

def _create_positioning_map(result: ResearchResult) -> bytes:
    """ポジショニングマップの画像を生成（ラベル重なり防止付き）"""
    fig, ax = plt.subplots(figsize=(7, 5))

    comp = result.competitor
    ax_label_x = comp.positioning_axis_x or "軸1"
    ax_label_y = comp.positioning_axis_y or "軸2"

    # ラベル位置の重なり防止用リスト
    placed_labels = []  # [(x, y)] 既配置ラベル座標

    def _get_offset(px, py):
        """既存ラベルと重ならないオフセットを計算"""
        offsets = [(8, 8), (-8, 8), (8, -12), (-8, -12), (12, 0), (-12, 0)]
        for ox, oy in offsets:
            conflict = False
            for lx, ly in placed_labels:
                if abs((px + ox / 10) - lx) < 0.8 and abs((py + oy / 10) - ly) < 0.6:
                    conflict = True
                    break
            if not conflict:
                placed_labels.append((px + ox / 10, py + oy / 10))
                return (ox, oy)
        placed_labels.append((px + 8 / 10, py + 8 / 10))
        return (8, 8)

    # 対象企業のプロット
    target_pos = getattr(comp, "_target_position", (5, 5))
    ax.scatter([target_pos[0]], [target_pos[1]], s=200, c="#E94560", zorder=5, marker="*")
    t_name = result.client_name[:10] + "…" if len(result.client_name) > 10 else result.client_name
    t_off = _get_offset(target_pos[0], target_pos[1])
    ax.annotate(t_name, (target_pos[0], target_pos[1]),
                fontsize=8, fontweight="bold", color="#E94560",
                xytext=t_off, textcoords="offset points")

    # 直接競合（四角マーカー）
    for c in comp.direct_competitors:
        ax.scatter([c.position_x], [c.position_y], s=80, c="#0F3460", zorder=4, marker="s")
        c_name = c.name[:10] + "…" if len(c.name) > 10 else c.name
        c_off = _get_offset(c.position_x, c.position_y)
        ax.annotate(c_name, (c.position_x, c.position_y),
                    fontsize=7, color="#333333",
                    xytext=c_off, textcoords="offset points")

    # 間接競合（ダイヤモンドマーカー）
    for c in comp.indirect_competitors:
        ax.scatter([c.position_x], [c.position_y], s=50, c="#999999", zorder=3, marker="D")
        c_name = c.name[:10] + "…" if len(c.name) > 10 else c.name
        c_off = _get_offset(c.position_x, c.position_y)
        ax.annotate(c_name, (c.position_x, c.position_y),
                    fontsize=7, color="#666666",
                    xytext=c_off, textcoords="offset points")

    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.set_xlabel(ax_label_x, fontsize=10)
    ax.set_ylabel(ax_label_y, fontsize=10)
    ax.set_title("ポジショニングマップ", fontsize=12, fontweight="bold")
    ax.grid(True, alpha=0.3)
    ax.axhline(y=5, color="#ccc", linestyle="--", linewidth=0.8)
    ax.axvline(x=5, color="#ccc", linestyle="--", linewidth=0.8)

    buf = io.BytesIO()
    fig.tight_layout()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _create_timeline(result: ResearchResult) -> bytes:
    """タイムライン画像を生成（最大8件、交互配置で重なり防止）"""
    events = result.company.timeline
    if not events:
        return b""

    fig, ax = plt.subplots(figsize=(10, 3.5))

    years = []
    labels = []
    for e in events:
        try:
            y = int(e.year[:4])
        except (ValueError, IndexError):
            continue
        years.append(y)
        # 説明文を15文字で切り詰め
        desc = e.description[:15] + "…" if len(e.description) > 15 else e.description
        labels.append(f"{e.year}\n{desc}")

    if not years:
        plt.close(fig)
        return b""

    # 最大8件に制限
    if len(years) > 8:
        years = years[:8]
        labels = labels[:8]

    y_pos = [0] * len(years)
    ax.scatter(years, y_pos, s=60, c="#0F3460", zorder=5)

    # 交互に上下に大きくずらして重なりを防止
    for i, (year, label) in enumerate(zip(years, labels)):
        if i % 2 == 0:
            offset_y = 30
            va = "bottom"
        else:
            offset_y = -30
            va = "top"
        ax.annotate(label, (year, 0), fontsize=6, ha="center", va=va,
                    xytext=(0, offset_y), textcoords="offset points",
                    arrowprops=dict(arrowstyle="-", color="#ccc", lw=0.5))

    ax.axhline(y=0, color="#0F3460", linewidth=2, alpha=0.5)
    ax.set_ylim(-1.5, 1.5)
    ax.set_yticks([])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.set_title("企業沿革", fontsize=12, fontweight="bold")

    buf = io.BytesIO()
    fig.tight_layout()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# --- スライド生成 ---

def _slide_cover(prs: Presentation, result: ResearchResult):
    """表紙スライド"""
    slide = _add_blank_slide(prs)

    # 背景色
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _hex_to_rgb(PptxStyle.PRIMARY)

    # タイトル
    left, top = Inches(1), Inches(2)
    width, height = Inches(10), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    _add_text(tf, f"{result.client_name}", size=36, bold=True, color=PptxStyle.TEXT_LIGHT, alignment=PP_ALIGN.LEFT)

    # サブタイトル
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(10), Inches(1))
    tf2 = txBox2.text_frame
    _add_text(tf2, "3C分析レポート", size=20, color="CCCCCC", alignment=PP_ALIGN.LEFT)

    # 日付
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(10), Inches(0.5))
    tf3 = txBox3.text_frame
    _add_text(tf3, f"分析実施日: {result.created_at}", size=12, color="999999", alignment=PP_ALIGN.LEFT)

    # 業界
    txBox4 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(10), Inches(0.5))
    tf4 = txBox4.text_frame
    _add_text(tf4, f"業界: {result.industry}", size=12, color="999999", alignment=PP_ALIGN.LEFT)


def _slide_executive_summary(prs: Presentation, result: ResearchResult):
    """エグゼクティブサマリー"""
    slide = _add_blank_slide(prs)

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    _add_text(tf, "Executive Summary", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

    # 区切り線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(11), Pt(3))
    _set_shape_bg(line, PptxStyle.HIGHLIGHT)
    line.line.fill.background()

    # 企業概要
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(11), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    overview = result.company.business_overview or f"{result.client_name}は{result.industry}業界の企業です。"
    _add_text(tf2, overview, size=PptxStyle.SIZE_BODY, color=PptxStyle.TEXT_DARK)

    # キーファインディング
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(11), Inches(0.5))
    tf3 = txBox3.text_frame
    _add_text(tf3, "Key Findings", size=PptxStyle.SIZE_HEADING, bold=True, color=PptxStyle.SECONDARY)

    y_offset = 2.8
    for i, finding in enumerate(result.key_findings[:5], 1):
        # ファインディングテキスト（番号付き）
        txBox_f = slide.shapes.add_textbox(Inches(0.7), Inches(y_offset), Inches(10.5), Inches(0.4))
        tf_f = txBox_f.text_frame
        tf_f.word_wrap = True
        _add_text(tf_f, f"{i}. {finding}", size=PptxStyle.SIZE_BODY, bold=True, color=PptxStyle.TEXT_DARK)

        y_offset += 0.55


def _slide_company(prs: Presentation, result: ResearchResult):
    """Company分析スライド（2〜3ページ）"""
    company = result.company

    # --- ページ1: 企業概要 ---
    slide1 = _add_blank_slide(prs)
    # タイトル
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    _add_text(tf, "Company - 企業概要", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(11), Pt(3))
    _set_shape_bg(line, PptxStyle.ACCENT)
    line.line.fill.background()

    # 企業情報ボックス
    info_items = [
        ("企業名", company.name),
        ("公式HP", company.official_url),
        ("理念・ビジョン", company.mission_vision),
        ("事業概要", company.business_overview),
        ("主要商品・サービス", company.products_services),
    ]

    y = 1.2
    for label, value in info_items:
        if not value:
            continue
        txBox_l = slide1.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2.5), Inches(0.3))
        tf_l = txBox_l.text_frame
        _add_text(tf_l, label, size=PptxStyle.SIZE_BODY, bold=True, color=PptxStyle.ACCENT)

        txBox_v = slide1.shapes.add_textbox(Inches(3.2), Inches(y), Inches(8.5), Inches(0.5))
        tf_v = txBox_v.text_frame
        tf_v.word_wrap = True
        # 長いテキストは短縮
        display_val = value[:200] + "..." if len(value) > 200 else value
        _add_text(tf_v, display_val, size=PptxStyle.SIZE_BODY, color=PptxStyle.TEXT_DARK)

        y += 0.7 if len(value) <= 80 else 1.0

    # --- ページ2: 沿革 + 最新動向 ---
    slide2 = _add_blank_slide(prs)
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    _add_text(tf, "Company - 沿革・最新動向", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

    line = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(11), Pt(3))
    _set_shape_bg(line, PptxStyle.ACCENT)
    line.line.fill.background()

    # タイムライン（python-pptx図形で構築、PPTX上で編集可能）
    events = company.timeline[:8]  # 最大8件
    if events:
        _draw_timeline_shapes(slide2, events)
        news_y = 3.6
    else:
        news_y = 1.2

    # 最新ニュース
    txBox_n = slide2.shapes.add_textbox(Inches(0.5), Inches(news_y), Inches(11), Inches(0.4))
    tf_n = txBox_n.text_frame
    _add_text(tf_n, "最新ニュース", size=PptxStyle.SIZE_HEADING, bold=True, color=PptxStyle.SECONDARY)

    # スライド下端までの残りスペースに応じてニュース件数を制限
    max_news_y = 6.8
    available = max_news_y - (news_y + 0.5)
    news_spacing = 0.85  # 各ニュース間のスペース
    max_news_count = max(1, int(available / news_spacing))
    display_news = company.recent_news[:min(4, max_news_count)]

    y = news_y + 0.5
    for news in display_news:
        if y + news_spacing > max_news_y:
            break
        txBox_item = slide2.shapes.add_textbox(Inches(0.7), Inches(y), Inches(10.5), Inches(0.7))
        tf_item = txBox_item.text_frame
        tf_item.word_wrap = True
        date_str = f"[{news.date}] " if news.date else ""
        title_text = news.title[:60] + "…" if len(news.title) > 60 else news.title
        _add_text(tf_item, f"{date_str}{title_text}", size=PptxStyle.SIZE_BODY, color=PptxStyle.TEXT_DARK)
        if news.summary:
            summary_text = news.summary[:80] + "…" if len(news.summary) > 80 else news.summary
            _add_text(tf_item, f"  {summary_text}", size=PptxStyle.SIZE_SMALL, color="666666")
        y += news_spacing

    # --- ページ3: SNS・ブランド評価 ---
    if company.sns_analysis or company.brand_momentum:
        slide3 = _add_blank_slide(prs)
        txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
        tf = txBox.text_frame
        _add_text(tf, "Company - ブランド評価・SNS分析", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

        line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(11), Pt(3))
        _set_shape_bg(line, PptxStyle.ACCENT)
        line.line.fill.background()

        # ブランドの勢い
        if company.brand_momentum:
            txBox_m = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(11), Inches(0.4))
            tf_m = txBox_m.text_frame
            _add_text(tf_m, "ブランドの勢い・熱量", size=PptxStyle.SIZE_HEADING, bold=True, color=PptxStyle.SECONDARY)

            txBox_mv = slide3.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(10.5), Inches(0.8))
            tf_mv = txBox_mv.text_frame
            tf_mv.word_wrap = True
            _add_text(tf_mv, company.brand_momentum, size=PptxStyle.SIZE_BODY, color=PptxStyle.TEXT_DARK)

        # SNS分析
        y = 2.8
        for sns in company.sns_analysis:
            txBox_s = slide3.shapes.add_textbox(Inches(0.5), Inches(y), Inches(11), Inches(0.35))
            tf_s = txBox_s.text_frame
            tone_color = {"ポジティブ": "27AE60", "ネガティブ": "E74C3C"}.get(sns.tone, PptxStyle.TEXT_DARK)
            _add_text(tf_s, f"{sns.platform}  [トーン: {sns.tone}]", size=PptxStyle.SIZE_BODY, bold=True, color=tone_color)

            txBox_sd = slide3.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(10.5), Inches(0.6))
            tf_sd = txBox_sd.text_frame
            tf_sd.word_wrap = True
            _add_text(tf_sd, sns.summary, size=PptxStyle.SIZE_BODY, color=PptxStyle.TEXT_DARK)
            if sns.key_topics:
                _add_text(tf_sd, f"主な話題: {', '.join(sns.key_topics)}", size=PptxStyle.SIZE_SMALL, color="666666")
            y += 1.0


def _slide_competitor(prs: Presentation, result: ResearchResult):
    """Competitor分析スライド（1〜2ページ）"""
    comp = result.competitor

    # --- ページ1: ポジショニングマップ ---
    slide1 = _add_blank_slide(prs)
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    _add_text(tf, "Competitor - ポジショニングマップ", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(11), Pt(3))
    _set_shape_bg(line, PptxStyle.HIGHLIGHT)
    line.line.fill.background()

    # ポジショニングマップ画像
    map_img = _create_positioning_map(result)
    img_stream = io.BytesIO(map_img)
    slide1.shapes.add_picture(img_stream, Inches(0.8), Inches(1.2), Inches(7.5), Inches(5))

    # 凡例（python-pptx shapeで構築 - PPTX上で編集可能）
    legend_x = 8.8
    legend_y = 1.3
    txBox_legend_title = slide1.shapes.add_textbox(
        Inches(legend_x), Inches(legend_y), Inches(3), Inches(0.3))
    _add_text(txBox_legend_title.text_frame, "凡例", size=9, bold=True, color=PptxStyle.SECONDARY)

    legend_items = [
        ("★", result.client_name[:12], "E94560"),
        ("■", "直接競合", "0F3460"),
        ("◆", "間接競合", "999999"),
    ]
    for li, (marker, label, color) in enumerate(legend_items):
        ly = legend_y + 0.35 + li * 0.3
        txBox_li = slide1.shapes.add_textbox(
            Inches(legend_x), Inches(ly), Inches(3), Inches(0.25))
        _add_text(txBox_li.text_frame, f"{marker} {label}", size=8, color=color)

    # 業界ポジション
    if comp.industry_position:
        pos_y = legend_y + 0.35 + len(legend_items) * 0.3 + 0.3
        txBox_pos = slide1.shapes.add_textbox(Inches(legend_x), Inches(pos_y), Inches(3), Inches(0.25))
        _add_text(txBox_pos.text_frame, "業界ポジション", size=9, bold=True, color=PptxStyle.SECONDARY)

        pos_text = comp.industry_position[:120] + "…" if len(comp.industry_position) > 120 else comp.industry_position
        txBox_posv = slide1.shapes.add_textbox(Inches(legend_x), Inches(pos_y + 0.3), Inches(3), Inches(2.5))
        txBox_posv.text_frame.word_wrap = True
        _add_text(txBox_posv.text_frame, pos_text, size=8, color=PptxStyle.TEXT_DARK)

    # --- ページ2: 競合比較表 ---
    all_competitors = comp.direct_competitors + comp.indirect_competitors
    if all_competitors:
        slide2 = _add_blank_slide(prs)
        txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
        tf = txBox.text_frame
        _add_text(tf, "Competitor - 競合比較表", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

        line = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(11), Pt(3))
        _set_shape_bg(line, PptxStyle.HIGHLIGHT)
        line.line.fill.background()

        # テーブル
        cols = 4  # 企業名, 概要, 強み, 差別化ポイント
        rows = len(all_competitors) + 1  # ヘッダー + データ行
        table_shape = slide2.shapes.add_table(rows, cols, Inches(0.3), Inches(1.2), Inches(11.4), Inches(5))
        table = table_shape.table

        # ヘッダー
        headers = ["企業名", "概要", "強み", "差別化ポイント"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_to_rgb(PptxStyle.PRIMARY)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = _hex_to_rgb(PptxStyle.TEXT_LIGHT)

        # データ行
        for row_idx, c in enumerate(all_competitors, 1):
            values = [c.name, c.description[:80], c.strengths[:80], c.differentiation[:80]]
            for col_idx, val in enumerate(values):
                cell = table.cell(row_idx, col_idx)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.color.rgb = _hex_to_rgb(PptxStyle.TEXT_DARK)
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_to_rgb(PptxStyle.BG_LIGHT)


def _slide_customer(prs: Presentation, result: ResearchResult):
    """Customer分析スライド（内容量に応じて1〜2ページ）"""
    customer = result.customer
    slide = _add_blank_slide(prs)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    _add_text(tf, "Customer - 市場・顧客分析", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(11), Pt(3))
    _set_shape_bg(line, "27AE60")
    line.line.fill.background()

    MAX_Y = 6.5  # スライド下端の安全マージン

    # 市場規模
    y = 1.2
    if customer.market_size:
        txBox_ms = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(11), Inches(0.3))
        tf_ms = txBox_ms.text_frame
        _add_text(tf_ms, "市場規模", size=PptxStyle.SIZE_HEADING, bold=True, color=PptxStyle.SECONDARY)

        market_text = customer.market_size[:200] + "…" if len(customer.market_size) > 200 else customer.market_size
        txBox_msv = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(10.5), Inches(0.5))
        tf_msv = txBox_msv.text_frame
        tf_msv.word_wrap = True
        _add_text(tf_msv, market_text, size=PptxStyle.SIZE_BODY, color=PptxStyle.TEXT_DARK)
        y += 1.1

    # 市場トレンド
    if customer.market_trend:
        txBox_mt = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(11), Inches(0.3))
        tf_mt = txBox_mt.text_frame
        _add_text(tf_mt, "市場トレンド", size=PptxStyle.SIZE_HEADING, bold=True, color=PptxStyle.SECONDARY)

        trend_text = customer.market_trend[:200] + "…" if len(customer.market_trend) > 200 else customer.market_trend
        txBox_mtv = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(10.5), Inches(0.7))
        tf_mtv = txBox_mtv.text_frame
        tf_mtv.word_wrap = True
        _add_text(tf_mtv, trend_text, size=PptxStyle.SIZE_BODY, color=PptxStyle.TEXT_DARK)
        y += 1.3

    # ターゲット顧客層
    if customer.target_segments or customer.target_description:
        # 残りスペースが足りなければ次のスライドへ
        if y + 1.0 > MAX_Y:
            slide = _add_blank_slide(prs)
            txBox_t2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
            tf_t2 = txBox_t2.text_frame
            _add_text(tf_t2, "Customer - 顧客分析（続き）", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)
            line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(11), Pt(3))
            _set_shape_bg(line2, "27AE60")
            line2.line.fill.background()
            y = 1.2

        txBox_tg = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(11), Inches(0.3))
        tf_tg = txBox_tg.text_frame
        _add_text(tf_tg, "ターゲット顧客層", size=PptxStyle.SIZE_HEADING, bold=True, color=PptxStyle.SECONDARY)

        txBox_tgv = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(10.5), Inches(0.7))
        tf_tgv = txBox_tgv.text_frame
        tf_tgv.word_wrap = True
        if customer.target_segments:
            segments_text = "・" + "\n・".join(s[:40] for s in customer.target_segments[:5])
            _add_text(tf_tgv, segments_text, size=PptxStyle.SIZE_BODY, color=PptxStyle.TEXT_DARK)
        if customer.target_description:
            desc_text = customer.target_description[:200] + "…" if len(customer.target_description) > 200 else customer.target_description
            _add_text(tf_tgv, desc_text, size=PptxStyle.SIZE_SMALL, color=PptxStyle.TEXT_DARK)
        y += 1.3

    # 類似事例
    if customer.similar_cases:
        # 残りスペースが足りなければ次のスライドへ
        if y + 1.0 > MAX_Y:
            slide = _add_blank_slide(prs)
            txBox_t3 = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
            tf_t3 = txBox_t3.text_frame
            _add_text(tf_t3, "Customer - 類似事例", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)
            line3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(11), Pt(3))
            _set_shape_bg(line3, "27AE60")
            line3.line.fill.background()
            y = 1.2

        txBox_sc = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(11), Inches(0.3))
        tf_sc = txBox_sc.text_frame
        _add_text(tf_sc, "類似事例・参考企業", size=PptxStyle.SIZE_HEADING, bold=True, color=PptxStyle.SECONDARY)
        y += 0.4

        for case in customer.similar_cases[:3]:
            if y + 1.2 > MAX_Y:
                break
            # 事例タイトル
            txBox_ct = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(10.5), Inches(0.3))
            _add_text(txBox_ct.text_frame, f"{case.company}（{case.industry}）", size=PptxStyle.SIZE_BODY, bold=True, color=PptxStyle.ACCENT)
            # 事例説明
            case_desc = case.description[:100] + "…" if len(case.description) > 100 else case.description
            txBox_cd = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.3), Inches(10.3), Inches(0.7))
            txBox_cd.text_frame.word_wrap = True
            _add_text(txBox_cd.text_frame, case_desc, size=9, color=PptxStyle.TEXT_DARK)
            if case.relevance:
                rel_text = case.relevance[:80] + "…" if len(case.relevance) > 80 else case.relevance
                _add_text(txBox_cd.text_frame, f"→ {rel_text}", size=8, color="666666")
            y += 1.2


def _slide_perspective(prs: Presentation, result: ResearchResult):
    """立場別ニーズ分析スライド（3カラム、テキストのみ - 装飾図形なし）"""
    perspective = result.perspective
    # データが空なら生成しない
    if not (perspective.executive.needs or perspective.frontline.needs or perspective.customer.needs):
        return

    slide = _add_blank_slide(prs)

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    _add_text(tf, "立場別ニーズ分析", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

    # 3カラムの定義
    columns = [
        {
            "title": "👔 経営者視点",
            "color": PptxStyle.ACCENT,
            "data": perspective.executive,
            "labels": ["必要なこと", "懸念事項", "成長機会"],
        },
        {
            "title": "🔧 現場視点",
            "color": "27AE60",
            "data": perspective.frontline,
            "labels": ["必要なこと", "懸念事項", "改善機会"],
        },
        {
            "title": "👤 顧客視点",
            "color": PptxStyle.HIGHLIGHT,
            "data": perspective.customer,
            "labels": ["求めていること", "不安・懸念", "理想の体験"],
        },
    ]

    col_width = 3.5
    col_gap = 0.25
    col_start_x = 0.4

    for col_idx, col in enumerate(columns):
        x = col_start_x + col_idx * (col_width + col_gap)
        data = col["data"]
        values = [data.needs, data.concerns, data.opportunities]

        # カラムヘッダー（テキストのみ、図形なし）
        txBox_h = slide.shapes.add_textbox(
            Inches(x), Inches(1.0), Inches(col_width), Inches(0.4),
        )
        _add_text(txBox_h.text_frame, col["title"], size=12, bold=True,
                  color=col["color"], alignment=PP_ALIGN.CENTER)

        # 各項目（needs / concerns / opportunities）
        item_y = 1.5
        item_spacing = 1.8

        for label, value in zip(col["labels"], values):
            # ラベル
            txBox_label = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(item_y), Inches(col_width - 0.2), Inches(0.25),
            )
            _add_text(txBox_label.text_frame, f"■ {label}", size=8, bold=True, color=col["color"])

            # 値テキスト（最大200文字、十分な高さを確保）
            display_val = value[:200] + "…" if len(value) > 200 else value
            txBox_val = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(item_y + 0.25), Inches(col_width - 0.2), Inches(1.45),
            )
            tf_val = txBox_val.text_frame
            tf_val.word_wrap = True
            _add_text(tf_val, display_val or "（情報なし）", size=8, color=PptxStyle.TEXT_DARK)

            item_y += item_spacing


def _slide_questions(prs: Presentation, result: ResearchResult):
    """考えるべき問いスライド（テキストのみ - 装飾図形なし）"""
    qa = result.questions
    if not qa or not qa.questions:
        return

    questions = qa.questions[:30]
    role_name = qa.role or "総合的なマーケティング担当者"
    total = len(questions)

    # 1スライドあたり最大12個で分割
    PER_SLIDE = 12
    page = 0

    while page * PER_SLIDE < total:
        start_idx = page * PER_SLIDE
        end_idx = min(start_idx + PER_SLIDE, total)
        chunk = questions[start_idx:end_idx]

        slide = _add_blank_slide(prs)

        # タイトル
        suffix = f"（{page + 1}）" if total > PER_SLIDE else ""
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
        tf = txBox.text_frame
        _add_text(tf, f"考えるべき問い{suffix}", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

        # ロール表示
        txBox_role = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(11), Inches(0.3))
        _add_text(txBox_role.text_frame, f"🎭 {role_name} の視点から", size=PptxStyle.SIZE_SMALL, bold=True, color="2980B9")

        # 問いリスト
        y = 1.3
        for i, q in enumerate(chunk, start_idx + 1):
            display_q = q[:120] + "…" if len(q) > 120 else q
            txBox_q = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(10.8), Inches(0.35))
            tf_q = txBox_q.text_frame
            tf_q.word_wrap = True
            _add_text(tf_q, f"{i}. {display_q}", size=9, color=PptxStyle.TEXT_DARK)
            y += 0.42

        page += 1


def _slide_appendix(prs: Presentation, result: ResearchResult):
    """付録 - 情報ソース一覧（テキストのみ - 装飾図形なし）"""
    slide = _add_blank_slide(prs)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    _add_text(tf, "付録 - 情報ソース一覧", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)

    txBox_note = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(11), Inches(0.3))
    _add_text(txBox_note.text_frame, f"分析実施日: {result.created_at}　|　情報ソース数: {len(result.sources)}件", size=PptxStyle.SIZE_SMALL, color="666666")

    y = 1.3
    for i, source in enumerate(result.sources[:20], 1):
        if y > 6.5:
            # 次のスライドへ
            slide = _add_blank_slide(prs)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.6))
            tf = txBox.text_frame
            _add_text(tf, "付録 - 情報ソース一覧（続き）", size=PptxStyle.SIZE_TITLE, bold=True, color=PptxStyle.PRIMARY)
            y = 1.0

        txBox_src = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(11), Inches(0.3))
        tf_src = txBox_src.text_frame
        title = source.title or source.url
        _add_text(tf_src, f"{i}. {title[:80]}", size=PptxStyle.SIZE_SMALL, color=PptxStyle.TEXT_DARK)
        _add_text(tf_src, f"   {source.url[:100]}", size=PptxStyle.SIZE_CAPTION, color="888888")
        y += 0.4


# --- メイン関数 ---

def build_pptx(result: ResearchResult, output_dir: str = "output") -> str:
    """3C分析結果からpptxファイルを生成

    Args:
        result: リサーチ結果
        output_dir: 出力ディレクトリ

    Returns:
        生成されたpptxファイルのパス
    """
    prs = Presentation()
    prs.slide_width = Emu(PptxStyle.SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(PptxStyle.SLIDE_HEIGHT_EMU)

    # スライド生成
    _slide_cover(prs, result)
    _slide_executive_summary(prs, result)
    _slide_company(prs, result)
    _slide_competitor(prs, result)
    _slide_customer(prs, result)
    _slide_perspective(prs, result)
    _slide_questions(prs, result)
    _slide_appendix(prs, result)

    # ファイル保存
    os.makedirs(output_dir, exist_ok=True)
    safe_name = result.client_name.replace("/", "_").replace("\\", "_")
    filename = f"3C分析_{safe_name}_{result.created_at.replace(':', '-').replace(' ', '_')}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)

    return filepath


def build_pptx_bytes(result: ResearchResult) -> bytes:
    """3C分析結果からpptxのバイトデータを生成（Streamlitダウンロード用）"""
    prs = Presentation()
    prs.slide_width = Emu(PptxStyle.SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(PptxStyle.SLIDE_HEIGHT_EMU)

    _slide_cover(prs, result)
    _slide_executive_summary(prs, result)
    _slide_company(prs, result)
    _slide_competitor(prs, result)
    _slide_customer(prs, result)
    _slide_perspective(prs, result)
    _slide_questions(prs, result)
    _slide_appendix(prs, result)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
